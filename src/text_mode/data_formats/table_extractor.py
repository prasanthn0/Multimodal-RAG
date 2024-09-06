from docx import Document
from pptx import Presentation
import pandas as pd
from src.config import settings
import tabula
from src.prompts import PredefinedPrompts
import traceback
from src.utils import get_summary


def extract_tables_from_pdf(file_path):
    df_tables = tabula.read_pdf(file_path, pages="all")
    return df_tables


def process_data(data_value, user_id, file, database, data_type):
    try:
        response = get_summary(data_value, PredefinedPrompts.summary_template)
        data={}
        if data_type=='table':
            data['content']=str(convert_keys_to_strings(data_value.to_dict()))
            data['transformed_content']=response
        elif data_type=='image':
            data['content']=str(data_value)
            data['transformed_content']=str(data_value)
        elif data_type=='text':
            data['content']=str(data_value['page_content'])
            data['transformed_content']=data_value['page_content'] 
   
        metadata={'file_name':file, 'user_id': user_id}
        database.store_vector(data, metadata)
    except Exception as e:
        print('Exception:', e)
        # Print traceback to understand the exception
        traceback.print_exc()

def convert_keys_to_strings(dictionary):
    """
    Recursively converts all numeric keys to strings in a dictionary.
    """
    if isinstance(dictionary, dict):
        return {str(key): convert_keys_to_strings(value) for key, value in dictionary.items()}
    elif isinstance(dictionary, list):
        return [convert_keys_to_strings(item) for item in dictionary]
    else:
        return dictionary
    
def read_tables_from_docx(docx_file):
    doc = Document(docx_file)
    tables_df = []
    for table in doc.tables:
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text)
            table_data.append(row_data)
        df = pd.DataFrame(table_data[1:], columns=table_data[0])  # Assuming first row is headers
        tables_df.append(df)
    return tables_df

def chunk_text_with_overlap(text):
    chunks = []
    chunk_size = settings.CHUNK_SIZE
    overlap =settings.CHUNK_OVERLAP
    start = 0
    end = chunk_size
    while start < len(text):
        chunks.append(text[start:end])
        start += chunk_size - overlap
        end = start + chunk_size
    return chunks

def extract_table_data_from_pptx(pptx_file):
    presentation = Presentation(pptx_file)
    table_data = []
    text_data = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                table_data.append([[cell.text for cell in row.cells] for row in table.rows])
            elif shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text_data.append(paragraph.text)
    dfs = []
    for table in table_data:
        df = pd.DataFrame(table[1:], columns=table[0])
        dfs.append(df)
    return dfs, text_data

def read_excel_data(file_path):
    # Read the Excel file
    xls = pd.ExcelFile(file_path)
    
    # Extract table data from each sheet
    sheet_names = xls.sheet_names
    table_data = []

    for sheet_name in sheet_names:
        df = xls.parse(sheet_name)
        table_data.append(df)
    return table_data
